#!/usr/bin/env python3
"""Convert a PPTX file to a TUM Typst presentation draft.

Usage:
    python3 pptx_to_typst.py input.pptx [output.typ]

Outputs a Typst draft to stdout (or the output file if given).
Prints a slide-by-slide summary to stderr.
"""

import sys
import re
from pathlib import Path


def _require_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Emu
        from pptx.enum.shapes import PP_PLACEHOLDER
        return Presentation, PP_PLACEHOLDER
    except ImportError:
        sys.exit("python-pptx is required: pip install python-pptx")


def _text_runs(shape):
    """Return plain text from a shape's text frame, stripping excess whitespace."""
    if not shape.has_text_frame:
        return ""
    lines = []
    for para in shape.text_frame.paragraphs:
        line = "".join(run.text for run in para.runs).strip()
        if line:
            lines.append(line)
    return "\n".join(lines)


def _slide_title(slide):
    """Extract the title text from a slide, or empty string."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_type == 14:  # PLACEHOLDER
            try:
                if shape.placeholder_format.idx == 0:
                    return _text_runs(shape)
            except Exception:
                pass
    # Fallback: largest text box
    candidates = [s for s in slide.shapes if s.has_text_frame]
    if candidates:
        biggest = max(candidates, key=lambda s: s.width * s.height)
        return _text_runs(biggest)
    return ""


def _slide_body_lines(slide):
    """Extract non-title body text lines."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    lines = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        is_title = False
        try:
            if shape.shape_type == 14 and shape.placeholder_format.idx == 0:
                is_title = True
        except Exception:
            pass
        if is_title:
            continue
        text = _text_runs(shape)
        if text:
            lines.extend(text.splitlines())
    return [l for l in lines if l.strip()]


def _has_significant_image(slide):
    """True if the slide contains an image placeholder or picture shape."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            return True
    return False


def _escape_typst(text):
    """Escape characters that have special meaning in Typst."""
    # Escape hash, underscore, asterisk, backtick, @, <, >
    for ch in ["#", "<", ">"]:
        text = text.replace(ch, "\\" + ch)
    return text


def _infer_slide_type(idx, slide, title, body_lines):
    """Heuristic: pick the best TUM slide type."""
    if idx == 0:
        return "title"
    if not body_lines and not _has_significant_image(slide):
        return "title-only"
    if _has_significant_image(slide) and len(body_lines) <= 2:
        return "image"
    return "content"


def _format_body(lines):
    """Turn a list of body lines into a Typst content block."""
    if not lines:
        return ""
    parts = []
    for line in lines:
        escaped = _escape_typst(line)
        # If it looks like a bullet point already (starts with -, •, *, numbers)
        if re.match(r"^[-•*\d]", line):
            parts.append(f"  {escaped}")
        else:
            parts.append(f"  {escaped}")
    return "\n".join(parts)


def convert(pptx_path: Path) -> str:
    Presentation, PP_PLACEHOLDER = _require_pptx()
    prs = Presentation(str(pptx_path))

    # Try to pull metadata from the first slide
    first_slide = prs.slides[0] if prs.slides else None
    cover_title = _slide_title(first_slide) if first_slide else "Presentation Title"
    cover_body = _slide_body_lines(first_slide) if first_slide else []

    # Guess author from first slide body (first non-empty line after title)
    author = cover_body[0] if cover_body else "Author"

    lines = [
        '#import "theme.typ": *',
        "",
        "#show: tum-theme.with(",
        f'  authors: ("{_escape_typst(author)}",),',
        f'  title: "{_escape_typst(cover_title)}",',
        '  footer-infos: (),',
        '  // school: "TUM School of ...",',
        '  // chair: "Lehrstuhl für ...",',
        ")",
        "",
        "// --- Title slide ---",
        "#title-slide()",
        "",
    ]

    for idx, slide in enumerate(prs.slides):
        title = _slide_title(slide)
        body_lines = _slide_body_lines(slide)
        slide_type = _infer_slide_type(idx, slide, title, body_lines)

        print(f"  Slide {idx+1}: [{slide_type:12s}] {title[:60]}", file=sys.stderr)

        if idx == 0:
            continue  # already emitted as #title-slide()

        if slide_type == "title-only":
            # Section divider — emit as a plain title slide with empty content
            lines.append(f"// --- Section divider ---")
            lines.append(f'#title-content-slide(title: "{_escape_typst(title)}")[')
            lines.append("]")
            lines.append("")
            continue

        if slide_type == "image":
            lines.append(f"// --- Image slide ---")
            lines.append(f'#title-image-slide(title: "{_escape_typst(title)}", image_path: "REPLACE_WITH_IMAGE_PATH")')
            lines.append("")
            continue

        # Default: title-content-slide
        lines.append(f"// --- Slide {idx+1} ---")
        lines.append(f'#title-content-slide(title: "{_escape_typst(title)}")[')
        body = _format_body(body_lines)
        if body:
            lines.append(body)
        else:
            lines.append("  // TODO: add content")
        lines.append("]")
        lines.append("")

    return "\n".join(lines)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        sys.exit(f"Usage: {sys.argv[0]} input.pptx [output.typ]")

    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        sys.exit(f"File not found: {pptx_path}")

    print(f"Converting {pptx_path} ...", file=sys.stderr)
    result = convert(pptx_path)

    if len(sys.argv) >= 3:
        out_path = Path(sys.argv[2])
        out_path.write_text(result, encoding="utf-8")
        print(f"Written to {out_path}", file=sys.stderr)
    else:
        print(result)
